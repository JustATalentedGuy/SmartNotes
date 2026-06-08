"""
extractor.py  —  revised
Extracts text and figures from PDF and PPTX files using pymupdf directly.

Key design decisions:
  • NO pymupdf4llm / magika dependency (avoids ONNX runtime errors)
  • Handles LaTeX article, Beamer slides, PPT-exported PDFs equally
  • Vector figure detection via page.get_drawings() bounding-box clustering
  • Beamer overlay deduplication via 3-gram text similarity
  • Heading detection via font-size analysis (median baseline)
  • Bullet normalisation (PDF bullet chars → markdown "- ")
  • PPTX: python-pptx with speaker notes and group-shape image extraction
"""

import os
import re
from pathlib import Path

SUPPORTED = {".pdf", ".pptx", ".ppt"}

# Typical Beamer footer band (bottom 8% of page height) — skip text there
_FOOTER_RATIO  = 0.92
# Thin header strip to skip (top 4% AND block height < 14 pt)
_HEADER_RATIO  = 0.04
_HEADER_MAX_H  = 14
# Minimum drawing rect dimension to be considered (filters degenerate lines)
_MIN_DRAW_DIM  = 2
# Gap (in PDF user-units) for merging nearby drawing bboxes into one figure
_MERGE_GAP     = 50
# Minimum merged figure area (sq. units) to be worth rendering
_MIN_FIG_AREA  = 3500    # ≈ 70×50 pt
# Fraction of figure area that can be text before we call it "text-dominated"
_TEXT_DOM_FRAC = 0.55
# A drawing rect spanning this fraction of page width OR height is a
# background/theme decoration, not a content figure — skip it
_BG_SPAN_FRAC  = 0.82
# A merged figure covering more than this fraction of page AREA is likely a
# full-slide background render or decorative frame — skip it
_MAX_FIG_AREA_FRAC = 0.58


# ─────────────────────────────────────────────────────────────────────────────
# Public entry point
# ─────────────────────────────────────────────────────────────────────────────

def extract_file(file_path: str, session_dir: str) -> dict:
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(file_path, session_dir)
    elif ext in (".pptx", ".ppt"):
        return _extract_pptx(file_path, session_dir)
    return _err(file_path, f"Unsupported type: {ext}")


# ─────────────────────────────────────────────────────────────────────────────
# PDF extractor — direct pymupdf (no pymupdf4llm)
# ─────────────────────────────────────────────────────────────────────────────

def _extract_pdf(pdf_path: str, session_dir: str) -> dict:
    """
    Extract text + figures from any PDF using pymupdf directly.

    Works for:
      • Standard LaTeX (article / report / book class)
      • Beamer presentation slides  (PPT-like formatting from LaTeX)
      • PPT/Keynote exported to PDF
      • Regular word-processor PDFs

    Figure detection:
      1. Raster images embedded in the PDF  (type-1 rawdict blocks)
      2. Vector graphic clusters             (get_drawings, bbox-merged)
    """
    try:
        import pymupdf
    except ImportError:
        return _err(pdf_path, "pymupdf missing — run: pip install pymupdf")

    img_dir = Path(session_dir) / "imgs"
    img_dir.mkdir(parents=True, exist_ok=True)
    source = Path(pdf_path).name
    stem   = Path(pdf_path).stem

    try:
        doc = pymupdf.open(pdf_path)
    except Exception as e:
        return _err(pdf_path, f"Cannot open: {e}")

    text_parts: list[str] = []
    all_images: list[dict] = []
    prev_text  = ""
    page_idx   = 0   # logical page counter (skips overlay duplicates)

    for pn in range(len(doc)):
        page     = doc[pn]
        page_md  = _page_to_markdown(page)

        # ── Deduplicate Beamer overlay pages ─────────────────────────────
        if _sim(page_md, prev_text) > 0.88:
            continue
        prev_text = page_md

        if page_md.strip():
            text_parts.append(page_md)

        # ── Figure extraction ─────────────────────────────────────────────
        figs = _extract_figures(page, img_dir, pn, stem, source)
        all_images.extend(figs)
        page_idx += 1

    doc.close()

    return {
        "source":   source,
        "type":     "pdf",
        "markdown": "\n\n---\n\n".join(text_parts),
        "images":   all_images,
    }


# ─────────────────────────────────────────────────────────────────────────────
# Page → Markdown  (text extraction with heading detection)
# ─────────────────────────────────────────────────────────────────────────────

# PDF-encoded bullet characters to normalise
_BULLET_RE = re.compile(r'^[•·▪▸►‣◦‐–—\uf0b7\uf0a7]+\s*')
# Lone page-number patterns (e.g. "3", "3 / 12", "Page 3")
_PAGE_NUM_RE = re.compile(r'^\s*(Page\s*)?\d{1,4}(\s*/\s*\d{1,4})?\s*$', re.IGNORECASE)

def _page_to_markdown(page) -> str:
    """Convert one PDF page to structured Markdown via font-size analysis."""
    try:
        raw = page.get_text("dict", sort=True)
    except Exception:
        return page.get_text("text").strip()

    blocks   = raw.get("blocks", [])
    page_h   = page.rect.height

    # ── Compute body font size (median across all meaningful spans) ──────
    sizes = [
        span["size"]
        for b in blocks if b.get("type") == 0
        for line in b.get("lines", [])
        for span in line.get("spans", [])
        if span.get("text", "").strip() and len(span["text"].strip()) > 2
    ]
    body_sz  = sorted(sizes)[len(sizes) // 2] if sizes else 10.0
    h2_thresh = body_sz * 1.35   # clearly a heading
    h3_thresh = body_sz * 1.15   # sub-heading

    out_lines: list[str] = []

    for block in blocks:
        if block.get("type") != 0:      # skip image blocks here
            continue

        bbox  = block.get("bbox", (0, 0, 0, 0))
        by0, by1 = bbox[1], bbox[3]
        blk_h = by1 - by0

        # Filter footer band
        if by0 > page_h * _FOOTER_RATIO:
            continue
        # Filter very thin header strip (but NOT slide titles which are tall)
        if by0 < page_h * _HEADER_RATIO and blk_h < _HEADER_MAX_H:
            continue

        for line in block.get("lines", []):
            parts, max_sz, is_bold = [], 0.0, False
            for span in line.get("spans", []):
                txt = span.get("text", "").strip()
                sz  = span.get("size", body_sz)
                flags = span.get("flags", 0)
                if txt:
                    parts.append(txt)
                    max_sz  = max(max_sz, sz)
                    is_bold = is_bold or bool(flags & 16)   # bit 4 = bold

            line_text = " ".join(parts).strip()
            if not line_text or len(line_text) < 2:
                continue
            # Drop lone page numbers / slide counters
            if _PAGE_NUM_RE.match(line_text):
                continue

            # ── Assign Markdown heading level ────────────────────────────
            if max_sz >= h2_thresh:
                out_lines.append(f"\n## {line_text}")
            elif max_sz >= h3_thresh:
                out_lines.append(f"\n### {line_text}")
            else:
                # Normalise PDF bullet characters
                cleaned = _BULLET_RE.sub("- ", line_text)
                # Wrap bold spans (heuristic: whole line was bold at body size)
                if is_bold and max_sz < h3_thresh:
                    cleaned = f"**{cleaned}**"
                out_lines.append(cleaned)

    return "\n".join(out_lines).strip()


# ─────────────────────────────────────────────────────────────────────────────
# Figure extraction  (raster + vector)
# ─────────────────────────────────────────────────────────────────────────────

def _extract_figures(page, img_dir: Path, page_num: int, stem: str, source: str) -> list:
    """
    Detect and render figure regions on a PDF page.

    Strategy A — Raster images:   type-1 blocks in rawdict
    Strategy B — Vector clusters: get_drawings(), bbox-merged with _MERGE_GAP
    Both strategies feed into a shared merge + filter pass.
    """
    import pymupdf

    candidate_rects: list = []     # pymupdf.Rect objects

    # ── A: raster image blocks ────────────────────────────────────────────
    try:
        raw = page.get_text("rawdict", sort=True)
        for block in raw.get("blocks", []):
            if block.get("type") == 1:
                r = pymupdf.Rect(block["bbox"])
                if r.width > 30 and r.height > 20:
                    candidate_rects.append(r)
    except Exception:
        pass

    # Also try get_images (catches images not in rawdict in some PDFs)
    try:
        for img_info in page.get_images(full=True):
            xref = img_info[0]
            try:
                for r in page.get_image_rects(xref):
                    fr = pymupdf.Rect(r)
                    if fr.width > 30 and fr.height > 20:
                        candidate_rects.append(fr)
            except Exception:
                pass
    except Exception:
        pass

    # ── B: vector drawings ───────────────────────────────────────────────
    try:
        pw, ph = page.rect.width, page.rect.height
        draw_rects = []
        for d in page.get_drawings():
            r = d.get("rect")
            if r is None:
                continue
            dr = pymupdf.Rect(r)
            # Skip degenerate zero-dimension paths (lines, dots)
            if dr.width <= _MIN_DRAW_DIM or dr.height <= _MIN_DRAW_DIM:
                continue
            # Skip background / theme decorations: any rect that spans most
            # of the page in either dimension is a slide background or border,
            # not a content figure (Beamer headers, full-page fills, etc.)
            if dr.width / pw > _BG_SPAN_FRAC or dr.height / ph > _BG_SPAN_FRAC:
                continue
            draw_rects.append(dr)

        # Cluster nearby drawing paths into figure groups
        merged_draws = _merge_rects(draw_rects, gap=_MERGE_GAP)
        for mr in merged_draws:
            area = mr.width * mr.height
            # Skip trivially small clusters
            if area < _MIN_FIG_AREA:
                continue
            # Skip merged regions that still cover too much of the page —
            # these are theme frames or the full content area of a slide
            if area / (pw * ph) > _MAX_FIG_AREA_FRAC:
                continue
            candidate_rects.append(mr)
    except Exception:
        pass

    if not candidate_rects:
        return []

    # ── Final merge of all candidate rects ───────────────────────────────
    all_merged = _merge_rects(candidate_rects, gap=12)

    # ── Get text blocks for context extraction and text-dom filtering ─────
    text_blocks: list[dict] = []
    try:
        td = page.get_text("dict", sort=True)
        text_blocks = [b for b in td.get("blocks", []) if b.get("type") == 0]
    except Exception:
        pass

    images: list[dict] = []
    seen: set[tuple] = set()

    for i, fig_rect in enumerate(all_merged):
        clipped = fig_rect & page.rect
        if clipped.is_empty:
            continue
        if clipped.width < 55 or clipped.height < 45:
            continue
        # Quantised dedup key
        key = (round(clipped.x0 / 8) * 8, round(clipped.y0 / 8) * 8,
               round(clipped.x1 / 8) * 8, round(clipped.y1 / 8) * 8)
        if key in seen:
            continue
        seen.add(key)

        if _text_dominated(clipped, text_blocks):
            continue

        # Final area-ratio guard: skip if the figure still covers too large a
        # fraction of the page (catches cases the earlier filter missed, e.g.
        # when raster image blocks and drawings merge into a full-page region)
        clipped_area = clipped.width * clipped.height
        page_area    = page.rect.width * page.rect.height
        if clipped_area / page_area > _MAX_FIG_AREA_FRAC:
            continue

        ctx = _nearby_text(clipped, text_blocks, page.rect)

        try:
            mat = pymupdf.Matrix(2.0, 2.0)     # 2× zoom → ~150 dpi at A4
            pix = page.get_pixmap(matrix=mat, clip=clipped, alpha=False,
                                  colorspace=pymupdf.csRGB)
            fname     = f"{stem}_p{page_num}_f{i}.png"
            img_path  = img_dir / fname
            pix.save(str(img_path))
            images.append({
                "path":        str(img_path.resolve()),
                "page":        page_num,
                "context":     ctx[:500],
                "source_file": source,
                "alt_text":    f"Figure from {source} p.{page_num + 1}",
            })
        except Exception:
            pass

    return images


# ─────────────────────────────────────────────────────────────────────────────
# Rectangle helpers
# ─────────────────────────────────────────────────────────────────────────────

def _merge_rects(rects: list, gap: int = 15) -> list:
    """Iteratively merge rectangles that are within `gap` units of each other."""
    if not rects:
        return []
    import pymupdf
    merged = [pymupdf.Rect(r) for r in rects]
    changed = True
    while changed:
        changed = False
        result, used = [], [False] * len(merged)
        for i in range(len(merged)):
            if used[i]:
                continue
            cur = pymupdf.Rect(merged[i])
            for j in range(i + 1, len(merged)):
                if used[j]:
                    continue
                exp = pymupdf.Rect(cur.x0 - gap, cur.y0 - gap,
                                   cur.x1 + gap, cur.y1 + gap)
                if not (exp & merged[j]).is_empty:
                    cur = pymupdf.Rect(
                        min(cur.x0, merged[j].x0), min(cur.y0, merged[j].y0),
                        max(cur.x1, merged[j].x1), max(cur.y1, merged[j].y1),
                    )
                    used[j] = True
                    changed = True
            result.append(cur)
        merged = result
    return merged


def _text_dominated(fig_rect, text_blocks: list) -> bool:
    """Return True when text covers more than _TEXT_DOM_FRAC of fig_rect area."""
    import pymupdf
    fig_area = fig_rect.width * fig_rect.height
    if fig_area < 1:
        return True
    cover = 0.0
    for b in text_blocks:
        inter = fig_rect & pymupdf.Rect(b["bbox"])
        if not inter.is_empty:
            cover += inter.width * inter.height
    return (cover / fig_area) > _TEXT_DOM_FRAC


def _nearby_text(fig_rect, text_blocks: list, page_rect, margin: int = 90) -> str:
    """Collect text from blocks within `margin` units above/below fig_rect."""
    import pymupdf
    zone = pymupdf.Rect(
        page_rect.x0,
        max(page_rect.y0, fig_rect.y0 - margin),
        page_rect.x1,
        min(page_rect.y1, fig_rect.y1 + margin),
    )
    texts = []
    for b in text_blocks:
        if not (zone & pymupdf.Rect(b["bbox"])).is_empty:
            for line in b.get("lines", []):
                for span in line.get("spans", []):
                    t = span.get("text", "").strip()
                    if t:
                        texts.append(t)
    return " ".join(texts)


# ─────────────────────────────────────────────────────────────────────────────
# Beamer overlay deduplication
# ─────────────────────────────────────────────────────────────────────────────

def _sim(a: str, b: str) -> float:
    """3-gram Jaccard similarity, fast approximate duplicate check."""
    if not a or not b or len(a) < 3 or len(b) < 3:
        return 0.0
    ga = {a[i:i+3] for i in range(len(a) - 2)}
    gb = {b[i:i+3] for i in range(len(b) - 2)}
    denom = len(ga | gb)
    return len(ga & gb) / denom if denom else 0.0


# ─────────────────────────────────────────────────────────────────────────────
# PPTX extractor  (unchanged logic, minor clean-up)
# ─────────────────────────────────────────────────────────────────────────────

def _extract_pptx(pptx_path: str, session_dir: str) -> dict:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return _err(pptx_path, "python-pptx missing — run: pip install python-pptx")

    img_dir = Path(session_dir) / "imgs"
    img_dir.mkdir(parents=True, exist_ok=True)
    source = Path(pptx_path).name

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return _err(pptx_path, f"Cannot open PPTX: {e}")

    text_parts, all_images = [], []

    for si, slide in enumerate(prs.slides):
        title_text, body_lines, notes_text = "", [], ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass

        for shape in slide.shapes:
            if shape.has_text_frame:
                is_title = (
                    hasattr(shape, "placeholder_format")
                    and shape.placeholder_format is not None
                    and shape.placeholder_format.idx == 0
                )
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if not line:
                        continue
                    if is_title and not title_text:
                        title_text = line
                    else:
                        lvl = getattr(para, "level", 0)
                        body_lines.append("  " * lvl + "- " + line)

            for stype in (MSO_SHAPE_TYPE.PICTURE,):
                if shape.shape_type == stype:
                    _save_pptx_image(shape, si, len(all_images),
                                     img_dir, source, title_text,
                                     body_lines, all_images)

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    for gs in shape.shapes:
                        if gs.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            _save_pptx_image(gs, si, len(all_images),
                                             img_dir, source, title_text,
                                             body_lines, all_images)
                except Exception:
                    pass

        slide_md = []
        if title_text:
            slide_md.append(f"## {title_text}")
        slide_md.extend(body_lines)
        if notes_text:
            slide_md.append(f"\n> 📝 **Lecturer notes:** {notes_text}")
        if slide_md:
            text_parts.append("\n".join(slide_md))

    return {
        "source":   source,
        "type":     "pptx",
        "markdown": "\n\n---\n\n".join(text_parts),
        "images":   all_images,
    }


def _save_pptx_image(shape, slide_idx: int, img_count: int,
                     img_dir: Path, source: str,
                     title: str, body_lines: list,
                     all_images: list) -> None:
    try:
        blob = shape.image.blob
        ext  = (shape.image.ext or "png").lower().lstrip(".")
        if ext not in ("png", "jpg", "jpeg", "gif", "bmp", "webp"):
            ext = "png"
        fname    = f"pptx_{Path(source).stem}_s{slide_idx}_i{img_count}.{ext}"
        img_path = img_dir / fname
        img_path.write_bytes(blob)
        ctx = " | ".join(filter(None, [title] + [l.lstrip("- ") for l in body_lines[:3]]))
        all_images.append({
            "path":        str(img_path.resolve()),
            "page":        slide_idx,
            "context":     ctx[:500],
            "source_file": source,
            "alt_text":    f"Slide {slide_idx + 1}: {title or 'figure'}",
        })
    except Exception:
        pass


# ─────────────────────────────────────────────────────────────────────────────
# Shared error helper
# ─────────────────────────────────────────────────────────────────────────────

def _err(file_path: str, msg: str) -> dict:
    return {
        "source":   Path(file_path).name,
        "type":     "error",
        "markdown": f"[EXTRACTION ERROR — {Path(file_path).name}: {msg}]",
        "images":   [],
        "error":    msg,
    }
